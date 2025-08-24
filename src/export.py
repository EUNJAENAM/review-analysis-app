"""
리포트 내보내기 모듈
PDF와 PPTX 형식으로 리포트를 생성합니다.
"""

import logging
import os
import re
from pathlib import Path
from typing import Dict, Any, Optional
import pandas as pd

from .config import OUTPUT_DIR, FONT_PATH, REPORT_METADATA

logger = logging.getLogger(__name__)

try:
    import subprocess
    import platform
    WKHTMLTOPDF_AVAILABLE = True
    # wkhtmltopdf가 시스템에 설치되어 있는지 확인
    try:
        subprocess.run(['wkhtmltopdf', '--version'], 
                      capture_output=True, check=True)
    except (subprocess.CalledProcessError, FileNotFoundError):
        WKHTMLTOPDF_AVAILABLE = False
        logger.warning("wkhtmltopdf가 시스템에 설치되지 않았습니다.")
except ImportError:
    WKHTMLTOPDF_AVAILABLE = False
    logger.warning("subprocess 모듈을 사용할 수 없습니다. PDF 생성을 건너뜁니다.")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False
    logger.warning("python-pptx 라이브러리가 설치되지 않았습니다. PPTX 생성을 건너뜁니다.")


class ExportGenerator:
    """PDF와 PPTX 리포트 생성 클래스"""
    
    def __init__(self, html_file_path: str):
        """
        Args:
            html_file_path: HTML 리포트 파일 경로
        """
        self.html_file_path = Path(html_file_path)
        self.output_dir = Path(OUTPUT_DIR)
        
    def _enhance_html_for_pdf(self, html_path: Path) -> Path:
        """
        PDF 생성을 위해 HTML 파일의 색상을 강화합니다.
        
        Args:
            html_path: 원본 HTML 파일 경로
            
        Returns:
            강화된 HTML 파일 경로
        """
        try:
            # HTML 파일 읽기
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 색상 강화 패턴들
            enhancements = [
                # KPI 카드 색상 강화
                (r'<div class="kpi-card">', '<div class="kpi-card" style="background: #4682B4 !important; color: white !important; border: 2px solid #4682B4; padding: 15px !important; margin-bottom: 10px !important;">'),
                # 테이블 헤더 색상 강화
                (r'<th>', '<th style="background-color: #4682B4 !important; color: white !important; border: 1px solid #4682B4; font-size: 11px !important; padding: 8px !important;">'),
                # 우선순위 색상 강화
                (r'class="priority-high"', 'class="priority-high" style="color: #DC143C !important; font-weight: bold !important;"'),
                (r'class="priority-medium"', 'class="priority-medium" style="color: #FF8C00 !important; font-weight: bold !important;"'),
                (r'class="priority-low"', 'class="priority-low" style="color: #32CD32 !important; font-weight: bold !important;"'),
                # 섹션 헤더 색상 강화
                (r'<h2>', '<h2 style="color: #4682B4 !important; border-bottom: 2px solid #4682B4 !important; font-size: 1.3em !important; margin-bottom: 10px !important;">'),
                # 컨테이너 배경 강화
                (r'<div class="container">', '<div class="container" style="background-color: white !important; padding: 15px !important;">'),
                # 차트 컨테이너 크기 조정
                (r'<div class="chart-container">', '<div class="chart-container" style="margin: 15px 0 !important; page-break-inside: avoid !important;">'),
                # 차트 이미지 크기 조정
                (r'<img src="{{ plot_files.', '<img style="max-width: 90% !important; height: auto !important; max-height: 400px !important;" src="{{ plot_files.'),
                # 테이블 셀 크기 조정
                (r'<td>', '<td style="padding: 8px !important; font-size: 11px !important;">'),
                # 섹션 크기 조정
                (r'<div class="section">', '<div class="section" style="margin-bottom: 20px !important; padding: 15px !important; page-break-inside: avoid !important;">'),
            ]
            
            # 색상 강화 적용
            for pattern, replacement in enhancements:
                html_content = re.sub(pattern, replacement, html_content)
            
            # 강화된 HTML 파일 저장
            enhanced_html_path = html_path.parent / f"enhanced_{html_path.name}"
            with open(enhanced_html_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            return enhanced_html_path
            
        except Exception as e:
            logger.warning(f"HTML 색상 강화 실패: {e}")
            return html_path
        
    def generate_pdf(self, output_filename: Optional[str] = None) -> Optional[str]:
        """
        HTML 리포트를 PDF로 변환 (wkhtmltopdf 사용)
        
        Args:
            output_filename: 출력 파일명 (기본값: report.pdf)
            
        Returns:
            생성된 PDF 파일 경로 또는 None (실패 시)
        """
        if not WKHTMLTOPDF_AVAILABLE:
            logger.error("wkhtmltopdf가 설치되지 않았습니다.")
            logger.info("설치 방법:")
            logger.info("1. https://wkhtmltopdf.org/downloads.html 에서 다운로드")
            logger.info("2. 또는 chocolatey: choco install wkhtmltopdf")
            logger.info("3. 또는 scoop: scoop install wkhtmltopdf")
            return None
            
        if not self.html_file_path.exists():
            logger.error(f"HTML 파일을 찾을 수 없습니다: {self.html_file_path}")
            return None
            
        try:
            if output_filename is None:
                output_filename = "report.pdf"
                
            pdf_path = self.output_dir / "report" / output_filename
            
            # PDF 디렉토리 생성
            pdf_path.parent.mkdir(parents=True, exist_ok=True)
            
            # HTML 색상 강화
            enhanced_html_path = self._enhance_html_for_pdf(self.html_file_path)
            
            # wkhtmltopdf 명령어 실행 (색상 렌더링 개선 + 크기 조정)
            cmd = [
                'wkhtmltopdf',
                '--page-size', 'A4',
                '--orientation', 'Portrait',
                '--margin-top', '0.5in',
                '--margin-right', '0.5in',
                '--margin-bottom', '0.5in',
                '--margin-left', '0.5in',
                '--encoding', 'UTF-8',
                '--print-media-type',
                '--enable-local-file-access',
                '--disable-smart-shrinking',
                '--dpi', '200',
                '--image-quality', '85',
                '--enable-forms',
                '--javascript-delay', '1000',
                '--no-stop-slow-scripts',
                '--load-error-handling', 'ignore',
                '--load-media-error-handling', 'ignore',
                '--zoom', '0.9',
                '--minimum-font-size', '8',
                '--disable-javascript',
                str(enhanced_html_path),
                str(pdf_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode == 0:
                logger.info(f"PDF 리포트가 생성되었습니다: {pdf_path}")
                return str(pdf_path)
            else:
                logger.error(f"wkhtmltopdf 실행 실패: {result.stderr}")
                return None
            
        except Exception as e:
            logger.error(f"PDF 생성 중 오류 발생: {e}")
            return None
    
    def generate_pptx(self, analysis_results: Dict[str, Any], 
                     plot_files: Dict[str, str],
                     output_filename: Optional[str] = None) -> Optional[str]:
        """
        분석 결과를 PPTX 요약 슬라이드로 생성
        
        Args:
            analysis_results: 분석 결과 데이터
            plot_files: 생성된 그래프 파일 경로들
            output_filename: 출력 파일명 (기본값: summary.pptx)
            
        Returns:
            생성된 PPTX 파일 경로 또는 None (실패 시)
        """
        if not PYTHON_PPTX_AVAILABLE:
            logger.error("python-pptx 라이브러리가 설치되지 않았습니다.")
            return None
            
        try:
            if output_filename is None:
                output_filename = "summary.pptx"
                
            pptx_path = self.output_dir / "report" / output_filename
            
            # PPTX 디렉토리 생성
            pptx_path.parent.mkdir(parents=True, exist_ok=True)
            
            # 프레젠테이션 생성
            prs = Presentation()
            
            # 슬라이드 레이아웃 선택 (제목과 콘텐츠)
            slide_layout = prs.slide_layouts[1]  # 제목과 콘텐츠 레이아웃
            slide = prs.slides.add_slide(slide_layout)
            
            # 제목 설정
            title = slide.shapes.title
            title.text = f"{REPORT_METADATA['hotel_name']} 리뷰 분석 요약"
            title.text_frame.paragraphs[0].font.size = Pt(28)
            title.text_frame.paragraphs[0].font.bold = True
            
            # 콘텐츠 영역
            content = slide.placeholders[1]
            content_frame = content.text_frame
            
            # KPI 카드 정보 추가
            kpis = analysis_results.get('핵심성과지표', {})
            if kpis:
                content_frame.text = "📊 핵심 성과지표"
                p = content_frame.paragraphs[0]
                p.font.size = Pt(18)
                p.font.bold = True
                
                # KPI 데이터 추가
                kpi_text = f"• 총 리뷰 수: {kpis.get('총리뷰수', 'N/A'):,}개\n"
                kpi_text += f"• 평균 평점: {kpis.get('평균평점', 'N/A'):.1f}점\n"
                kpi_text += f"• 긍정 비율: {kpis.get('긍정비율', 'N/A'):.1f}%\n"
                kpi_text += f"• 부정 비율: {kpis.get('부정비율', 'N/A'):.1f}%\n"
                
                p = content_frame.add_paragraph()
                p.text = kpi_text
                p.font.size = Pt(14)
                
            # 주요 인사이트 추가
            insights = analysis_results.get('주요인사이트', [])
            if insights:
                p = content_frame.add_paragraph()
                p.text = "\n💡 주요 인사이트"
                p.font.size = Pt(18)
                p.font.bold = True
                
                insight_text = ""
                for i, insight in enumerate(insights[:3], 1):  # 상위 3개만
                    insight_text += f"{i}. {insight}\n"
                
                p = content_frame.add_paragraph()
                p.text = insight_text
                p.font.size = Pt(12)
            
            # 개선 우선순위 추가
            priorities = analysis_results.get('개선우선순위', [])
            if priorities:
                p = content_frame.add_paragraph()
                p.text = "\n🎯 개선 우선순위"
                p.font.size = Pt(18)
                p.font.bold = True
                
                priority_text = ""
                for i, priority in enumerate(priorities[:3], 1):  # 상위 3개만
                    priority_text += f"{i}. {priority}\n"
                
                p = content_frame.add_paragraph()
                p.text = priority_text
                p.font.size = Pt(12)
            
            # 전략적 제언 추가
            strategies = analysis_results.get('전략적제언', {})
            if strategies:
                p = content_frame.add_paragraph()
                p.text = "\n📈 전략적 제언"
                p.font.size = Pt(18)
                p.font.bold = True
                
                strategy_text = ""
                for period, strategy in strategies.items():
                    if strategy:
                        strategy_text += f"• {period}: {str(strategy)[:50]}...\n"
                
                p = content_frame.add_paragraph()
                p.text = strategy_text
                p.font.size = Pt(12)
            
            # 프레젠테이션 저장
            prs.save(str(pptx_path))
            
            logger.info(f"PPTX 요약 슬라이드가 생성되었습니다: {pptx_path}")
            return str(pptx_path)
            
        except Exception as e:
            logger.error(f"PPTX 생성 중 오류 발생: {e}")
            return None
    
    def generate_all_formats(self, analysis_results: Dict[str, Any], 
                           plot_files: Dict[str, str]) -> Dict[str, Optional[str]]:
        """
        모든 형식의 리포트 생성
        
        Args:
            analysis_results: 분석 결과 데이터
            plot_files: 생성된 그래프 파일 경로들
            
        Returns:
            생성된 파일 경로들의 딕셔너리
        """
        results = {}
        
        # PDF 생성
        pdf_path = self.generate_pdf()
        results['pdf'] = pdf_path
        
        # PPTX 생성
        pptx_path = self.generate_pptx(analysis_results, plot_files)
        results['pptx'] = pptx_path
        
        return results
